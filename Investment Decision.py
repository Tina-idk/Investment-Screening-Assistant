import streamlit as st
st.title("Investment Screening Assistant")
st.write("Upload company documents and let AI help evaluate whether the business meets your investment criteria!")
uploaded_file = st.file_uploader("Upload company profile, pitch deck, or business plan", type=["pdf", "pptx"])

# Text Extraction
def extract_text(uploaded_file):
    file_type = uploaded_file.name.split('.')[-1].lower()
    if file_type == "pdf":
        import fitz
        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        return "\n".join([page.get_text() for page in doc])
    elif file_type == "pptx":
        from pptx import Presentation
        prs = Presentation(uploaded_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    else:
        return "Unsupported file type."
        
def generate_multi_comparison_conclusion(records):
    companies_text = ""
    for i, record in enumerate(records, 1):
        companies_text += f"\nCompany {i} ({record['filename']}):\n{record['score']}\n"

    prompt = f"""
    Compare the following companies based on their investment scores.
    Each has been evaluated based on: Product, Team, Financial Assessment, Market Attractiveness, Exit Potential.

    {companies_text}

    Write a comparison summary (5â€“8 sentences) explaining:
    - Which companies perform better overall and why
    - Their respective strengths and weaknesses
    - If any company stands out as the best investment target
    - Any ties or close comparisons worth mentioning

    Be concise, neutral, and insightful.
    """
    response = model.generate_content(prompt, generation_config={"temperature": 0})
    return response.text
    
def parse_score_table_to_df(markdown_table_str):
    import pandas as pd
    from io import StringIO

    lines = markdown_table_str.strip().splitlines()
    table_lines = [line for line in lines if line.strip().startswith("|") and "---" not in line]
    if len(table_lines) < 2:
        return pd.DataFrame()

    csv_str = "\n".join([line.strip().strip("|").replace(" | ", ",") for line in table_lines[1:]])
    df = pd.read_csv(StringIO(csv_str), names=["Criteria", "Score", "Explanation"])
    return df

def extract_scores_only(markdown_table_str):
    df = parse_score_table_to_df(markdown_table_str)
    return df.set_index("Criteria")["Score"]


# AI Analysis
import google.generativeai as genai
genai.configure(api_key=st.secrets["API"])
def analyze_with_ai(text):
    model = genai.GenerativeModel(model_name="gemini-1.5-flash")
    def split_text(text, chunk_size=3000):
        return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]
    summaries = []
    for i, chunk in enumerate(split_text(text)):
        prompt = f"""
        This is part {i + 1} of a company document.
        Try to be concise and accurate. Do not omit any information that may affect investment decisions.
        Summarize the key points of this section only:

        {chunk}
        """
        response = model.generate_content(prompt, generation_config={"temperature": 0})
        summaries.append(response.text)
    final_input = "\n\n".join(summaries)

    intro_prompt = f"""
    You are an investment analyst. Read the company content below and provide:
    
    - Company Description
    - Founders / Team
    - Main product/service
    - Key Features
    - Historical Financial Performance
    - Market Overview
    - Future Plan
    - Key Risks and Miligants
    
    Content:
    {final_input}
    """
    intro_response = model.generate_content(intro_prompt, generation_config={"temperature": 0}).text
    
    SEIS_promt = f"""
    Please assess whether the following company meets the eligibility criteria for the UK Enterprise Investment Scheme (EIS).
    Evaluate each criterion separately and respond with:
    âœ… Yes / âŒ No / âš ï¸ Uncertain
    
    Criterion:
    UK-registered company
    Age â‰¤ 2 years
    Assets â‰¤ Â£350, 000
    Employees â‰¤ 25
    Risk-to-capital condition
    Business operates in a qualifying trade (not financial, real estate, energy generation, etc.)
    Funds will be used for eligible business purposes (not repaying debt or buying companies)
        
    Return as markdown table, with columns Criterion, Status, Explanation.
    Conclude whether the company is likely to qualify for EIS.
    
    Content:
    {final_input}
    """
    SEIS_response = model.generate_content(SEIS_promt, generation_config={"temperature": 0}).text

    EIS_promt = f"""
    Please assess whether the following company meets the eligibility criteria for the UK Enterprise Investment Scheme (EIS).
    
    Evaluate each criterion separately and respond with:
    âœ… Yes / âŒ No / âš ï¸ Uncertain
    
    Criterion:
    UK-registered company
    Unlisted (or listed on AIM)
    Age â‰¤ 7 years
    Assets â‰¤ Â£15 million
    Employees â‰¤ 250
    Risk-to-capital condition
    Business operates in a qualifying trade (not financial, real estate, energy generation, etc.)
    Funds will be used for eligible business purposes (not repaying debt or buying companies)
    
    Return as markdown table, with columns Criterion, Status, Explanation.
    Conclude whether the company is likely to qualify for EIS.
    
    Content:
    {final_input}
    """
    EIS_response = model.generate_content(EIS_promt, generation_config={"temperature": 0}).text

    score_prompt = f"""
    Evaluate the company based on these criteria and corresponding standards of evaluation after the colon:
    Annual Revenue: Â£0 â€“ Â£50K -> score 1ï½ž2 (Pre-revenue, idea-stage, may be SEIS)
                    Â£50K â€“ Â£200K -> score 3ï½ž4 (MVP with some customers)
                    Â£200K â€“ Â£500K -> score 5ï½ž6 (Early monetization, but still proving model)
                    Â£500K â€“ Â£1M -> score 7ï½ž8 (Growing revenue, signs of traction)
                    Â£1M+ -> score 9ï½ž10 (Commercial validation, possible scaling)
    YoY Growth: Calculate the YoY Growth rate based on: (Year2 - Year1) / Year1 (use latest data)
                <0% â†’ Score 1ï½ž2
                0â€“50% â†’ Score 3ï½ž4
                50â€“100% â†’ Score 5ï½ž6
                100â€“200% â†’ Score 7ï½ž8
                >200% â†’ Score 9ï½ž10
    Founder & Team Assessment: 
        No relevant experience or unverifiable team â†’ Score 1ï½ž2
        Solo founder or minimal relevant background â†’ Score 3ï½ž4
        Small team with some relevant background â†’ Score 5ï½ž6
        Strong, diverse team with good experience â†’ Score 7ï½ž8
        Proven, highly credible founding team with strong track record â†’ Score 9ï½ž10
    Products/Services:
        No product yet / purely conceptual â†’ Score 1ï½ž2
        MVP exists but little to no validation â†’ Score 3ï½ž4
        Functional with early users/pilots â†’ Score 5ï½ž6
        PMF indicators, clear adoption â†’ Score 7ï½ž8
        Proven product with strong differentiation and growth â†’ Score 9ï½ž10
    Exit Potential:
        No exit plan, no precedent â†’ Score 1ï½ž2
        Vague or unrealistic exit ideas â†’ Score 3ï½ž4
        Some exit potential, unclear path â†’ Score 5ï½ž6
        Clear scenarios, comparables or likely buyers exist â†’ Score 7ï½ž8
        Strong exit potential with multiple paths and founder experience â†’ Score 9ï½ž10
        
    For each:
    - Score out of 10
    - 1-sentence explanation (with numeric data)

    Return as markdown table. Show the final score out of 50. Conclude with investment recommendation.

    Content:
    {final_input}
    """
    score_response = model.generate_content(score_prompt, generation_config={"temperature": 0}).text

    return intro_response, SEIS_response, EIS_response, score_response

# Output
# Only show "Save Result" if AI analysis has been run
if "analysis_done" not in st.session_state:
    st.session_state["analysis_done"] = False

if "records" not in st.session_state:
    st.session_state["records"] = []
    
if uploaded_file:
    content = extract_text(uploaded_file)
    if content:
        st.text_area("Document Preview", content[:2000])
        if st.button("Analyze with AI"):
            with st.spinner("Analyzing..."):
                intro_response, SEIS_response, EIS_response, score_response = analyze_with_ai(content)
                st.session_state["intro_response"] = intro_response
                st.session_state["score_response"] = score_response
                st.session_state["analysis_done"] = True  # âœ… Mark as done
            st.subheader("Company Overview")
            st.write(intro_response)
            st.subheader("SEIS Judgement")
            st.markdown(SEIS_response)
            st.subheader("EIS Judgement")
            st.markdown(EIS_response)
            st.subheader("Investment Criteria Evaluation")
            st.markdown(score_response)
    else:
        st.error("Unsupported file type or empty content.")

MAX_RECORDS = 10

# Only allow saving if analysis is completed
if st.session_state.get("analysis_done", False):
    if st.button("Save Result"):
        st.session_state["records"].append({
            "filename": uploaded_file.name,
            "overview": st.session_state["intro_response"],
            "score": st.session_state["score_response"],
        })
        if len(st.session_state["records"]) > MAX_RECORDS:
            st.session_state["records"] = st.session_state["records"][-MAX_RECORDS:]
        st.success("Result saved!")

if "records" in st.session_state and len(st.session_state["records"]) > 0:
    options = [f"{i+1}. {record['filename']}" for i, record in enumerate(st.session_state["records"])]
    default_selection = options[:2] if len(options) >= 2 else options

    selected = st.multiselect(
        "Select companies to compare (up to 10)",
        options,
        default=default_selection,
        max_selections=10
    )

    if len(selected) >= 2:
        indices = [int(s.split(".")[0]) - 1 for s in selected]
        records = [st.session_state["records"][i] for i in indices]
        score_dict = {}
        for record in records:
            score_dict[record["filename"]] = extract_scores_only(record["score"]) 
        df = pd.DataFrame(score_dict)
        df.index.name = "Criteria"
        st.markdown("### Score Table Comparison")
        st.dataframe(df, use_container_width=True)
        summary = generate_multi_comparison_conclusion(records)
        st.subheader("AI Summary Conclusion")
        st.write(summary)
   
        full_tables = []
        for record in records:
            df_full = parse_score_table_to_df(record["score"])
            df_full["Source File"] = record["filename"] 
            full_tables.append(df_full)
        if full_tables:
            combined_df = pd.concat(full_tables, ignore_index=True)
            st.markdown("### ðŸ§¾ Full Score Table with Explanation")
            st.dataframe(combined_df, use_container_width=True)
            csv = combined_df.to_csv(index=False).encode("utf-8")
            st.download_button(
                label="ðŸ“¥ Download full table as CSV",
                data=csv,
                file_name="full_scores_with_explanations.csv",
                mime="text/csv"
            )

    elif len(selected) == 1:
        idx = int(selected[0].split(".")[0]) - 1
        record = st.session_state["records"][idx]
        st.markdown(f"###  {record['filename']}")
        st.markdown(record["score"])

else:
    st.info("Please upload at least one company file.")


if st.button("Clear All Saved Analyses"):
    st.session_state["records"] = []
    st.success("All records have been cleared!")
